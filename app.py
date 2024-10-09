from flask import Flask, render_template, request, redirect, url_for, send_file
# Importa os módulos necessários do Flask:
# - Flask: para criar o aplicativo web.
# - render_template: para renderizar templates HTML.
# - request: para acessar dados da solicitação HTTP.
# - redirect: para redirecionar o usuário para uma URL diferente.
# - url_for: para gerar URLs para funções de view.
# - send_file: para enviar arquivos como resposta HTTP.

import mysql.connector
# Importa mysql.connector para conectar e interagir com o banco de dados MySQL.

import pandas as pd
# Importa pandas para manipulação e análise de dados, especialmente útil para DataFrames e leitura/escrita de arquivos.

from docx import Document
# Importa Document do módulo docx para criar e manipular documentos no formato DOCX (Microsoft Word).

from pptx import Presentation
from pptx.util import Inches
# Importa Presentation e Inches do módulo pptx para criar e manipular apresentações no formato PPTX (Microsoft PowerPoint),
# e para definir tamanhos e posições dos elementos dos slides.

import io
# Importa io para manipulação de fluxo de dados em memória, útil para ler e escrever arquivos em bytes.

import hashlib
# Importa hashlib para criar hashes criptográficos, usados para gerar códigos curtos únicos.

from datetime import datetime
# Importa datetime para manipulação de datas e horários.

import plotly.express as px
import plotly.io as pio
# Importa plotly.express para criar gráficos interativos de forma simplificada.
# Importa plotly.io para exportar gráficos como imagens.

app = Flask(__name__)
# Cria uma instância do aplicativo Flask.

# Configuração do banco de dados
db_config = {
    'user': 'root',                # Nome de usuário para o banco de dados MySQL.
    'password': '',                # Senha para o banco de dados MySQL (vazia neste caso).
    'host': 'localhost',           # Endereço do servidor MySQL.
    'database': 'tca'              # Nome do banco de dados.
}

db = mysql.connector.connect(**db_config)
# Estabelece uma conexão com o banco de dados MySQL usando as configurações fornecidas.

def extract_short_code(url):
    """Extrai as 4 primeiras letras do nome da URL, ignorando 'www' e 'https'."""
    url = url.replace('https://', '')  # Remove o prefixo 'https://' se existir.
    url = url.replace('http://', '')   # Remove o prefixo 'http://' se existir.
    url = url.replace('www.', '')       # Remove o prefixo 'www.' se existir.
    short_code = url[:4].upper()        # Extrai as 4 primeiras letras da URL e as transforma em maiúsculas.
    return f"WEB-{short_code}"          # Retorna o código curto no formato "Web-XXXX".

def generate_unique_short_code(original_url):
    """Gera um código curto único para a URL fornecida."""
    short_code = extract_short_code(original_url)  # Obtém o código curto inicial da URL original.
    while True:
        cursor = db.cursor()
        # Verifica se o código curto já existe no banco de dados.
        cursor.execute("SELECT COUNT(*) FROM urls WHERE short_code = %s", (short_code,))
        if cursor.fetchone()[0] == 0:
            cursor.close()
            return short_code  # Se o código não existir, retorna o código curto.
        # Se o código já existir, gera um novo código curto.
        short_code = extract_short_code(original_url) + str(datetime.now().microsecond)
        cursor.close()

@app.route("/", methods=["GET", "POST"])
def index():
    """Página principal para encurtar URLs."""
    if request.method == "POST":
        # Recebe a URL original do formulário enviado via POST.
        original_url = request.form["original_url"]
        # Gera um código curto único para a URL.
        short_code = generate_unique_short_code(original_url)

        cursor = db.cursor()
        # Insere a URL original, o código curto, a data de criação e a contagem de cliques inicial no banco de dados.
        cursor.execute("INSERT INTO urls (original_url, short_code, created_at, click_count) VALUES (%s, %s, %s, %s)",
                       (original_url, short_code, datetime.now(), 0))
        db.commit()  # Confirma a transação no banco de dados.
        cursor.close()

        # Gera a URL curta completa para redirecionar o usuário.
        short_url = url_for("redirect_url", short_code=short_code, _external=True)
        return render_template("index.html", short_url=short_url)  # Renderiza a página principal com a URL curta gerada.

    return render_template("index.html")  # Renderiza a página principal para solicitações GET.

@app.route("/<short_code>")
def redirect_url(short_code):
    """Redireciona para a URL original com base no código curto fornecido."""
    cursor = db.cursor(dictionary=True)
    # Recupera a URL original associada ao código curto.
    cursor.execute("SELECT original_url FROM urls WHERE short_code = %s", (short_code,))
    url_data = cursor.fetchone()

    if url_data:
        client_ip = request.remote_addr  # Obtém o IP do cliente.
        # Atualiza o registro no banco de dados com a data e IP do último clique, e incrementa a contagem de cliques.
        cursor.execute("""
            UPDATE urls 
            SET last_click_at = %s, last_click_ip = %s, click_count = click_count + 1 
            WHERE short_code = %s
        """, (datetime.now(), client_ip, short_code))
        db.commit()  # Confirma a transação no banco de dados.
        cursor.close()
        return redirect(url_data['original_url'])  # Redireciona o usuário para a URL original.
    else:
        cursor.close()
        return "URL não encontrada", 404  # Retorna um erro 404 se o código curto não for encontrado.

@app.route("/urls")
def show_urls():
    """Exibe todas as URLs encurtadas e suas informações."""
    cursor = db.cursor(dictionary=True)
    # Recupera todas as URLs do banco de dados.
    cursor.execute("SELECT * FROM urls")
    urls = cursor.fetchall()
    cursor.close()
    return render_template("urls.html", urls=urls)  # Renderiza a página com a lista de todas as URLs encurtadas.

@app.route("/charts")
def charts():
    """Exibe um gráfico com a contagem de cliques para cada URL."""
    cursor = db.cursor(dictionary=True)
    # Recupera dados das URLs do banco de dados.
    cursor.execute("SELECT short_code, original_url, click_count, created_at, last_click_at, last_click_ip FROM urls")
    urls = cursor.fetchall()
    cursor.close()

    # Prepara os dados para o gráfico.
    df = pd.DataFrame(urls)
    fig = px.bar(df, x='short_code', y='click_count', labels={'short_code': 'Short Code', 'click_count': 'Clicks'}, title='URL Clicks')

    # Converte o gráfico para HTML.
    chart_html = fig.to_html(full_html=False)

    return render_template("charts.html", url_data=urls, chart_html=chart_html)  # Renderiza a página com o gráfico das URLs e os dados.

@app.route("/download/<file_type>")
def download_report(file_type):
    """Gera e baixa um relatório no formato especificado (pptx, docx, xlsx)."""
    cursor = db.cursor(dictionary=True)
    # Recupera dados das URLs do banco de dados.
    cursor.execute("SELECT short_code, original_url, click_count, created_at, last_click_at, last_click_ip FROM urls")
    urls = cursor.fetchall()
    cursor.close()

    if file_type == 'pptx':
        # Cria uma apresentação PowerPoint.
        prs = Presentation()

        # Gera o gráfico e adiciona ao PowerPoint.
        df = pd.DataFrame(urls)
        fig = px.bar(df, x='short_code', y='click_count', labels={'short_code': 'Short Code', 'click_count': 'Clicks'}, title='URL Clicks')
        img_bytes = pio.to_image(fig, format='png')

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "URL Clicks"

        # Adiciona o gráfico ao slide.
        image_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(image_stream, Inches(1), Inches(1.5), width=Inches(8))

        # Adiciona informações das URLs a cada slide.
        for url in urls:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = f"Short Code: {url['short_code']}"

            content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
            text_frame = content.text_frame
            p = text_frame.add_paragraph()
            p.text = (
                f"Original URL: {url['original_url']}\n"
                f"Clicks: {url['click_count']}\n"
                f"Created At: {url['created_at'].strftime('%Y-%m-%d %H:%M:%S')}\n"
                f"Last Click: {url['last_click_at'].strftime('%Y-%m-%d %H:%M:%S') if url['last_click_at'] else 'N/A'}\n"
                f"Last Click IP: {url['last_click_ip'] or 'N/A'}"
            )

        # Salva a apresentação em memória.
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return send_file(output, as_attachment=True, download_name='report.pptx')

    elif file_type == 'xlsx':
        # Cria um arquivo Excel.
        df = pd.DataFrame(urls)
        output = io.BytesIO()
        with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
            # Escreve os dados no Excel.
            df.to_excel(writer, sheet_name='URLs', index=False)
            # Não é necessário chamar writer.save(); isso acontece automaticamente ao sair do bloco with.
        output.seek(0)

        return send_file(output, as_attachment=True, download_name='report.xlsx')


    elif file_type == 'docx':
        # Cria um documento Word.
        doc = Document()
        doc.add_heading('URL Report', 0)

        # Adiciona informações de cada URL ao documento.
        for url in urls:
            doc.add_heading(f"Short Code: {url['short_code']}", level=1)
            doc.add_paragraph(
                f"Original URL: {url['original_url']}\n"
                f"Clicks: {url['click_count']}\n"
                f"Created At: {url['created_at'].strftime('%Y-%m-%d %H:%M:%S')}\n"
                f"Last Click: {url['last_click_at'].strftime('%Y-%m-%d %H:%M:%S') if url['last_click_at'] else 'N/A'}\n"
                f"Last Click IP: {url['last_click_ip'] or 'N/A'}"
            )

        # Salva o documento em memória.
        output = io.BytesIO()
        doc.save(output)
        output.seek(0)

        return send_file(output, as_attachment=True, download_name='report.docx')

    else:
        return "Invalid file type", 400  # Retorna um erro 400 se o tipo de arquivo for inválido.

# Rota para exibir e registrar links
@app.route('/register_url', methods=['GET', 'POST'])
def register_url():
    if request.method == 'POST':
        name = request.form['name']
        link_type = request.form['link_type']
        url = request.form['url']

        # Conectar ao banco de dados e registrar o link
        conn = mysql.connector.connect(**db_config)
        cursor = conn.cursor()
        cursor.execute("INSERT INTO links (name, link_type, url) VALUES (%s, %s, %s)", (name, link_type, url))
        conn.commit()
        cursor.close()
        conn.close()

        return redirect('/register_url')  # Redirecionar após o cadastro

    # Recuperar os links cadastrados com cliques e últimos IPs
    conn = mysql.connector.connect(**db_config)
    cursor = conn.cursor(dictionary=True)
    cursor.execute("""
        SELECT 
            l.id AS link_id,
            l.name,
            l.link_type,
            l.url,
            l.created_at,
            COALESCE(COUNT(c.id), 0) AS clicks,
            COALESCE(MAX(c.ip_address), 'N/A') AS last_ip
        FROM 
            links l
        LEFT JOIN 
            clicks_links c ON l.id = c.url_id
        GROUP BY 
            l.id
    """)
    links = cursor.fetchall()
    cursor.close()
    conn.close()

    return render_template('search.html', links=links)  # Renderizar a página com os links cadastrados


# Função para exibir os links cadastrados
@app.route("/search", methods=["GET"])
def search():
    query = request.args.get("q")
    search_path = request.args.get("path")

    # Código para verificar e processar a pesquisa
    if not search_path:
        return render_template("search.html", error="Erro: O parâmetro de caminho não foi fornecido."), 400

    # Processar o caminho se for uma URL
    if is_url(search_path):
        result = search_online(search_path)
        return render_template("search.html", query=query, result=result)

    # Se for um diretório local
    else:
        results = search_local_files(search_path, query)
        return render_template("search.html", query=query, results=results)

        # Rota para registrar cliques
# Rota para registrar cliques
@app.route('/click_link/<int:link_id>', methods=['GET'])
def click_link(link_id):
    # Capturar o endereço IP do usuário
    ip_address = request.remote_addr

    # Conectar ao banco de dados e registrar o clique
    conn = mysql.connector.connect(**db_config)
    cursor = conn.cursor()

    # Inserir registro de clique
    cursor.execute("INSERT INTO clicks_links (url_id, click_time, ip_address) VALUES (%s, NOW(), %s)", (link_id, ip_address))
    conn.commit()

    # Obter a URL original para redirecionar
    cursor.execute("SELECT url FROM links WHERE id = %s", (link_id,))
    original_url = cursor.fetchone()
    
    cursor.close()
    conn.close()

    if original_url:
        return redirect(original_url[0])  # Redirecionar para a URL original

    return "URL não encontrada", 404  # Caso a URL não exista


 
if __name__ == "__main__":
    app.run(debug=True)

